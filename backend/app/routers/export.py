import io
import csv
from datetime import date as _date
from typing import Optional

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session

from app.deps import get_db
from app.models import Process, Application
from app.models.bia import BiaAssessment

router = APIRouter(prefix="/export", tags=["export"])

# ── Constants ─────────────────────────────────────────────────────────────────

SCORE_LABELS = {1: "Vitaal", 2: "Hoog", 3: "Midden", 4: "Laag", 5: "Minimaal"}
SCORE_BIA = {1: "Catastrofaal", 2: "Kritiek / Zeer ernstig", 3: "Gemiddeld", 4: "Gering", 5: "Verwaarloosbaar"}

# Tijdslabels per score, afgeleid van de antwoordschalen van B1 (uitvalduur → RTO)
# en B2 (dataverlies → RPO) in de BIA-vragenlijst.
RTO_LABELS = {1: "Enkele uren", 2: "Maximaal 8 uur", 3: "Maximaal 2 werkdagen", 4: "Maximaal 1 week", 5: "Meer dan een week"}
RPO_LABELS = {1: "Enkele uren", 2: "4–8 uur", 3: "8–24 uur", 4: "Maximaal 24 uur", 5: "Een week of meer"}

B_QUESTIONS = [
    ("b1", "B1 – Maximale uitvalduur (RTO)"),
    ("b2", "B2 – Maximale dataverlies (RPO)"),
    ("b3", "B3 – Hersteltijd achterstanden (WRT)"),
    ("b4", "B4 – Totale maximale uitvalduur (MTPD)"),
]
I_QUESTIONS = [("i1", "I1 – Impact onjuiste informatie")]
V_QUESTIONS = [("v1", "V1 – Impact ongeautoriseerde inzage")]

ALL_SECTIONS: dict[str, list[str]] = {
    "dashboard":         ["kpi", "biv_verdeling", "kritieke_processen", "review", "acties"],
    "processes":         ["basis", "details", "biv", "rto_rpo", "applicaties", "datums"],
    "applications":      ["basis", "details", "review", "processen"],
    "bia":               ["algemeen", "beschikbaarheid", "integriteit", "vertrouwelijkheid", "eindscores"],
    "business-context":  ["canvas", "wettelijk", "privacy", "continuiteit"],
    "ketenarchitectuur": ["processen", "applicaties", "koppelingen"],
}

SECTION_LABELS: dict[str, str] = {
    "kpi": "KPI Overzicht",
    "biv_verdeling": "BIV-verdeling",
    "kritieke_processen": "Kritieke Processen",
    "review": "Review Monitoring",
    "acties": "Prioritaire Acties",
    "basis": "Basisgegevens",
    "details": "Details",
    "biv": "BIV-scores",
    "rto_rpo": "RTO / RPO",
    "applicaties": "Gekoppelde Applicaties",
    "datums": "Datums",
    "processen": "Processen",
    "algemeen": "Algemene Informatie",
    "beschikbaarheid": "Beschikbaarheid (B)",
    "integriteit": "Integriteit (I)",
    "vertrouwelijkheid": "Vertrouwelijkheid (V)",
    "eindscores": "Eindscores BIV",
    "canvas": "Canvas Blokken",
    "wettelijk": "Wettelijk & Stakeholders",
    "privacy": "Privacy",
    "continuiteit": "Continuïteit",
    "koppelingen": "Koppelingen",
}

# ── Helpers ───────────────────────────────────────────────────────────────────

def _parse_sections(sections_param: Optional[str], module: str) -> list[str]:
    if not sections_param:
        return ALL_SECTIONS.get(module, [])
    return [s.strip() for s in sections_param.split(",") if s.strip()]


def _has_rto_rpo(p: Process) -> bool:
    """Zelfde definitie als het dashboard: RTO (b1) en RPO (b2) in de BIA ingevuld."""
    return p.bia is not None and p.bia.b1_score is not None and p.bia.b2_score is not None


def _rto_str(p: Process) -> str:
    if p.bia is None or p.bia.b1_score is None:
        return ""
    return RTO_LABELS.get(p.bia.b1_score, str(p.bia.b1_score))


def _rpo_str(p: Process) -> str:
    if p.bia is None or p.bia.b2_score is None:
        return ""
    return RPO_LABELS.get(p.bia.b2_score, str(p.bia.b2_score))


def _check_completeness(p: Process) -> tuple[bool, list[str]]:
    missing: list[str] = []
    if not p.description:       missing.append("Beschrijving")
    if not p.objective:         missing.append("Doelstelling")
    if not p.owner:             missing.append("Eigenaar")
    if not p.department:        missing.append("Afdeling")
    if not p.last_assessment_date: missing.append("Laatste beoordelingsdatum")
    if p.is_critical and not p.critical_reason: missing.append("Reden kritiek")
    if not p.applications:      missing.append("Gekoppelde applicaties")
    if p.bia is None:           missing.append("BIA / BIV")
    if not _has_rto_rpo(p):     missing.append("RTO / RPO")
    if p.business_context is None: missing.append("Business context")
    return len(missing) == 0, missing


def _score_str(score: Optional[int], labels: dict) -> str:
    if score is None:
        return ""
    return f"{score} – {labels.get(score, str(score))}"


# ── XLSX helpers ──────────────────────────────────────────────────────────────

# openpyxl weigert XML-illegale control-tekens (bv. verticale tab \x0b uit
# Word geplakt); vervang ze door een newline zodat de export niet crasht.
_XL_ILLEGAL_RE = __import__("re").compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")


def _xl_clean(value):
    if isinstance(value, str):
        return _XL_ILLEGAL_RE.sub("\n", value)
    return value


def _xl_append(ws, row: list):
    ws.append([_xl_clean(v) for v in row])


def _xl_header(ws, headers: list[str], fill_color: str = "003366"):
    from openpyxl.styles import Font, PatternFill, Alignment
    hfill = PatternFill("solid", fgColor=fill_color)
    hfont = Font(color="FFFFFF", bold=True)
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.fill = hfill
        cell.font = hfont
        cell.alignment = Alignment(horizontal="center")


def _xl_autowidth(ws):
    for col in ws.columns:
        max_len = max((len(str(c.value or "")) for c in col), default=0)
        ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)


def _xl_section_heading(ws, row: int, text: str, ncols: int):
    from openpyxl.styles import Font, PatternFill, Alignment
    ws.cell(row=row, column=1, value=text).font = Font(bold=True, size=12, color="FFFFFF")
    ws.cell(row=row, column=1).fill = PatternFill("solid", fgColor="1A5276")
    ws.cell(row=row, column=1).alignment = Alignment(horizontal="left")
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=max(ncols, 2))


# ── DOCX helpers ──────────────────────────────────────────────────────────────

def _doc_heading(doc, text: str, level: int = 1):
    doc.add_heading(text, level=level)


def _doc_table(doc, headers: list[str], rows: list[list]):
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Table Grid"
    # Header row
    hdr = table.rows[0]
    for i, h in enumerate(headers):
        cell = hdr.cells[i]
        cell.text = h
        run = cell.paragraphs[0].runs[0]
        run.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell._tc.get_or_add_tcPr().append(
            _tc_shading("003366")
        )
    # Data rows
    for r_idx, row in enumerate(rows, 1):
        for c_idx, val in enumerate(row):
            table.rows[r_idx].cells[c_idx].text = str(val) if val is not None else ""
    doc.add_paragraph()


def _tc_shading(color_hex: str):
    from docx.oxml import OxmlElement
    shd = OxmlElement("w:shd")
    shd.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill", color_hex)
    shd.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}color", "auto")
    shd.set("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val", "clear")
    return shd


def _doc_kv(doc, items: list[tuple[str, str]]):
    from docx.shared import Pt
    table = doc.add_table(rows=len(items), cols=2)
    table.style = "Table Grid"
    for i, (k, v) in enumerate(items):
        table.rows[i].cells[0].text = k
        table.rows[i].cells[0].paragraphs[0].runs[0].bold = True
        table.rows[i].cells[1].text = str(v) if v is not None else ""
    doc.add_paragraph()


# ── PPTX helpers ──────────────────────────────────────────────────────────────

def _prs_title_slide(prs, title: str, subtitle: str = ""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle and slide.placeholders[1] is not None:
        slide.placeholders[1].text = subtitle
    return slide


def _prs_content_slide(prs, title: str, content: str):
    from pptx.util import Pt
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = content
    tf.word_wrap = True
    return slide


_PPTX_ROWS_PER_SLIDE = 12  # meer rijen loopt van de slide af


def _prs_table_slide(prs, title: str, headers: list[str], rows: list[list]):
    """Rendert een tabel, verdeeld over meerdere slides bij veel rijen."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    layout = prs.slide_layouts[5]  # title only

    if not rows:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
        txBox.text_frame.text = "Geen gegevens beschikbaar."
        return slide

    chunks = [rows[i:i + _PPTX_ROWS_PER_SLIDE] for i in range(0, len(rows), _PPTX_ROWS_PER_SLIDE)]
    slide = None
    for idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title if idx == 0 else f"{title} (vervolg {idx + 1})"

        nrows = len(chunk) + 1
        ncols = len(headers)
        left, top = Inches(0.3), Inches(1.5)
        width, height = Inches(9.4), Inches(0.35 * nrows + 0.4)

        table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table

        # Header row
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)
            tf = cell.text_frame.paragraphs[0]
            tf.runs[0].font.bold = True
            tf.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tf.runs[0].font.size = Pt(9)

        # Data rows
        for r, row in enumerate(chunk, 1):
            for c, val in enumerate(row):
                cell = table.cell(r, c)
                cell.text = str(val) if val is not None else ""
                tf = cell.text_frame.paragraphs[0]
                if tf.runs:
                    tf.runs[0].font.size = Pt(8)

    return slide


# ── DATA GETTERS ──────────────────────────────────────────────────────────────

def _get_processes(db: Session) -> list[Process]:
    return db.query(Process).order_by(Process.code).all()


def _get_applications(db: Session) -> list[Application]:
    return db.query(Application).order_by(Application.code).all()


def _get_review_data(db: Session) -> dict:
    today = _date.today()
    try:
        cutoff = today.replace(year=today.year - 1)
    except ValueError:
        cutoff = _date(today.year - 1, 2, 28)

    def on_time(d) -> bool:
        return d is not None and d >= cutoff

    processes = db.query(Process).all()
    applications = db.query(Application).all()
    total_proc = len(processes)
    total_apps = len(applications)

    def _pct(done, total):
        return round(done / total * 100) if total > 0 else 0

    proc_done = sum(1 for p in processes if on_time(p.last_assessment_date))
    bia_done  = sum(1 for p in processes if p.bia is not None and on_time(p.bia.interview_date))
    bc_done   = sum(1 for p in processes if p.business_context is not None and on_time(p.business_context.review_date))
    apps_done = sum(1 for a in applications if on_time(a.review_date))

    return {
        "Processen": (proc_done, total_proc, _pct(proc_done, total_proc)),
        "Applicaties": (apps_done, total_apps, _pct(apps_done, total_apps)),
        "BIA assessments": (bia_done, total_proc, _pct(bia_done, total_proc)),
        "Business Context": (bc_done, total_proc, _pct(bc_done, total_proc)),
    }


def _get_priority_actions(processes: list[Process]) -> list[dict]:
    def _is_high_risk(p: Process) -> bool:
        if p.bia is None:
            return False
        for attr in ("availability_score", "integrity_score", "confidentiality_score"):
            s = getattr(p.bia, attr, None)
            if s is not None and s <= 2:
                return True
        return False

    actions = []
    for p in processes:
        _, missing = _check_completeness(p)
        if not missing:
            continue
        if p.is_critical and p.bia is None:
            priority, reason = "Kritiek", "Informatie ontbreekt"
        elif _is_high_risk(p) and not _has_rto_rpo(p):
            priority, reason = "Hoog", "Hoog risico: geen RTO/RPO gedefinieerd"
        elif p.is_critical:
            priority, reason = "Hoog", "Kritisch proces — onvolledig gedocumenteerd"
        else:
            priority, reason = "Middel", f"{len(missing)} veld(en) ontbreken"
        actions.append({
            "code": p.code, "naam": p.name,
            "kritiek": "Ja" if p.is_critical else "Nee",
            "prioriteit": priority, "reden": reason,
            "ontbrekend": ", ".join(missing),
        })
    order = {"Kritiek": 0, "Hoog": 1, "Middel": 2}
    return sorted(actions, key=lambda a: order.get(a["prioriteit"], 9))


# ══════════════════════════════════════════════════════════════════════════════
#  XLSX BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def _build_xlsx_dashboard(db: Session, sections: list[str]) -> io.BytesIO:
    from openpyxl import Workbook
    wb = Workbook()
    wb.remove(wb.active)  # remove default sheet
    processes = _get_processes(db)
    total = len(processes)

    def _cov(done):
        return round(done / total * 100) if total > 0 else 0

    if "kpi" in sections:
        ws = wb.create_sheet("KPI Overzicht")
        _xl_header(ws, ["Indicator", "Waarde", "Percentage"])
        complete = sum(1 for p in processes if _check_completeness(p)[0])
        attention = sum(1 for p in processes if 0 < len(_check_completeness(p)[1]) <= 3)
        incomplete = sum(1 for p in processes if len(_check_completeness(p)[1]) > 3)
        bia_cov = sum(1 for p in processes if p.bia is not None)
        bc_cov = sum(1 for p in processes if p.business_context is not None)
        apps_cov = sum(1 for p in processes if p.applications)
        rows = [
            ("Totaal processen", total, ""),
            ("Kritieke processen", sum(1 for p in processes if p.is_critical), f"{_cov(sum(1 for p in processes if p.is_critical))}%"),
            ("Volledig compleet", complete, f"{_cov(complete)}%"),
            ("Aandacht vereist (1–3 ontbrekend)", attention, f"{_cov(attention)}%"),
            ("Incompleet (>3 ontbrekend)", incomplete, f"{_cov(incomplete)}%"),
            ("BIA / BIV gedekt", bia_cov, f"{_cov(bia_cov)}%"),
            ("Business Context gedekt", bc_cov, f"{_cov(bc_cov)}%"),
            ("Applicaties gekoppeld", apps_cov, f"{_cov(apps_cov)}%"),
        ]
        for i, row in enumerate(rows, 2):
            for j, val in enumerate(row, 1):
                ws.cell(row=i, column=j, value=_xl_clean(val))
        _xl_autowidth(ws)

    if "biv_verdeling" in sections:
        ws = wb.create_sheet("BIV-verdeling")
        _xl_header(ws, ["Dimensie", "Vitaal (1)", "Hoog (2)", "Midden (3)", "Laag (4)", "Minimaal (5)", "Niet beoordeeld"])
        for dim, attr, label in [
            ("Beschikbaarheid", "availability_score", "B"),
            ("Integriteit", "integrity_score", "I"),
            ("Vertrouwelijkheid", "confidentiality_score", "V"),
        ]:
            counts = {1: 0, 2: 0, 3: 0, 4: 0, 5: 0}
            not_assessed = 0
            for p in processes:
                if p.bia is None:
                    not_assessed += 1
                else:
                    score = getattr(p.bia, attr, None)
                    if score in counts:
                        counts[score] += 1
                    else:
                        not_assessed += 1
            _xl_append(ws, [dim, counts[1], counts[2], counts[3], counts[4], counts[5], not_assessed])
        _xl_autowidth(ws)

    if "kritieke_processen" in sections:
        ws = wb.create_sheet("Kritieke Processen")
        _xl_header(ws, ["Code", "Naam", "B-score", "I-score", "V-score", "Heeft BIA", "Heeft RTO/RPO", "Ontbrekende velden"])
        for p in processes:
            if not p.is_critical:
                continue
            _, missing = _check_completeness(p)
            _xl_append(ws, [
                p.code, p.name,
                _score_str(p.bia.availability_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.integrity_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.confidentiality_score if p.bia else None, SCORE_LABELS),
                "Ja" if p.bia else "Nee",
                "Ja" if _has_rto_rpo(p) else "Nee",
                ", ".join(missing),
            ])
        _xl_autowidth(ws)

    if "review" in sections:
        ws = wb.create_sheet("Review Monitoring")
        _xl_header(ws, ["Categorie", "Op tijd", "Totaal", "Percentage"])
        rev = _get_review_data(db)
        for cat, (done, total_n, pct) in rev.items():
            _xl_append(ws, [cat, done, total_n, f"{pct}%"])
        _xl_autowidth(ws)

    if "acties" in sections:
        ws = wb.create_sheet("Prioritaire Acties")
        _xl_header(ws, ["Code", "Naam", "Kritiek", "Prioriteit", "Reden", "Ontbrekende velden"])
        for a in _get_priority_actions(processes):
            _xl_append(ws, [a["code"], a["naam"], a["kritiek"], a["prioriteit"], a["reden"], a["ontbrekend"]])
        _xl_autowidth(ws)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


def _build_xlsx_processes(db: Session, sections: list[str]) -> io.BytesIO:
    from openpyxl import Workbook
    processes = _get_processes(db)
    wb = Workbook()
    wb.remove(wb.active)

    # Determine which columns to include based on sections
    if "basis" in sections or "details" in sections or "biv" in sections or "rto_rpo" in sections or "applicaties" in sections or "datums" in sections:
        headers = []
        if "basis" in sections:
            headers += ["Code", "Naam", "Eigenaar", "Afdeling", "Kritiek", "Reden kritiek"]
        if "details" in sections:
            headers += ["Beschrijving", "Doelstelling", "Notities"]
        if "biv" in sections:
            headers += ["B-score", "I-score", "V-score"]
        if "rto_rpo" in sections:
            headers += ["RTO", "RPO"]
        if "applicaties" in sections:
            headers += ["# Applicaties", "Applicaties"]
        if "datums" in sections:
            headers += ["Laatste beoordeling", "Aangemaakt op"]

        ws = wb.create_sheet("Processen")
        _xl_header(ws, headers)
        for p in processes:
            row = []
            if "basis" in sections:
                row += [p.code, p.name, p.owner or "", p.department or "", "Ja" if p.is_critical else "Nee", p.critical_reason or ""]
            if "details" in sections:
                row += [p.description or "", p.objective or "", p.notes or ""]
            if "biv" in sections:
                row += [
                    _score_str(p.bia.availability_score if p.bia else None, SCORE_LABELS),
                    _score_str(p.bia.integrity_score if p.bia else None, SCORE_LABELS),
                    _score_str(p.bia.confidentiality_score if p.bia else None, SCORE_LABELS),
                ]
            if "rto_rpo" in sections:
                row += [_rto_str(p), _rpo_str(p)]
            if "applicaties" in sections:
                row += [len(p.applications), ", ".join(a.name for a in p.applications)]
            if "datums" in sections:
                row += [
                    str(p.last_assessment_date) if p.last_assessment_date else "",
                    str(p.created_at.date()) if p.created_at else "",
                ]
            _xl_append(ws, row)
        _xl_autowidth(ws)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


def _build_xlsx_applications(db: Session, sections: list[str]) -> io.BytesIO:
    from openpyxl import Workbook
    applications = _get_applications(db)
    wb = Workbook()
    wb.remove(wb.active)

    headers = []
    if "basis" in sections:
        headers += ["Code", "Naam", "Business owner", "Technisch owner"]
    if "details" in sections:
        headers += ["Beschrijving", "Notities"]
    if "review" in sections:
        headers += ["Review datum"]
    if "processen" in sections:
        headers += ["# Processen", "Processen"]

    if headers:
        ws = wb.create_sheet("Applicaties")
        _xl_header(ws, headers)
        for a in applications:
            row = []
            if "basis" in sections:
                row += [a.code, a.name, a.business_owner or "", a.technical_owner or ""]
            if "details" in sections:
                row += [a.description or "", a.notes or ""]
            if "review" in sections:
                row += [str(a.review_date) if a.review_date else ""]
            if "processen" in sections:
                row += [len(a.processes), ", ".join(p.name for p in a.processes)]
            _xl_append(ws, row)
        _xl_autowidth(ws)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


def _build_xlsx_bia(db: Session, sections: list[str]) -> io.BytesIO:
    from openpyxl import Workbook
    processes = _get_processes(db)
    wb = Workbook()
    wb.remove(wb.active)

    if "eindscores" in sections:
        ws = wb.create_sheet("Eindscores BIV")
        _xl_header(ws, ["Proces code", "Proces naam", "B-score", "I-score", "V-score"])
        for p in processes:
            _xl_append(ws, [
                p.code, p.name,
                _score_str(p.bia.availability_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.integrity_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.confidentiality_score if p.bia else None, SCORE_LABELS),
            ])
        _xl_autowidth(ws)

    if "algemeen" in sections:
        ws = wb.create_sheet("Algemeen")
        _xl_header(ws, ["Proces", "Interviewer", "Interview datum", "Beschrijving", "Ketenafhankelijkheden", "Afwijking motivatie", "Notities"])
        for p in processes:
            if p.bia is None:
                continue
            _xl_append(ws, [
                f"{p.code} – {p.name}",
                p.bia.interviewer_name or "",
                str(p.bia.interview_date) if p.bia.interview_date else "",
                p.bia.general_description or "",
                p.bia.chain_dependencies or "",
                p.bia.owner_deviation_motivation or "",
                p.bia.notes or "",
            ])
        _xl_autowidth(ws)

    if "beschikbaarheid" in sections:
        ws = wb.create_sheet("Beschikbaarheid (B)")
        q_headers = ["Proces"] + [label for _, label in B_QUESTIONS] + ["Eindscore B"]
        _xl_header(ws, q_headers)
        for p in processes:
            if p.bia is None:
                continue
            row = [f"{p.code} – {p.name}"]
            for key, _ in B_QUESTIONS:
                score = getattr(p.bia, f"{key}_score", None)
                row.append(_score_str(score, SCORE_BIA))
            row.append(_score_str(p.bia.availability_score, SCORE_LABELS))
            _xl_append(ws, row)
        _xl_autowidth(ws)

    if "integriteit" in sections:
        ws = wb.create_sheet("Integriteit (I)")
        q_headers = ["Proces"] + [label for _, label in I_QUESTIONS] + ["Eindscore I"]
        _xl_header(ws, q_headers)
        for p in processes:
            if p.bia is None:
                continue
            row = [f"{p.code} – {p.name}"]
            for key, _ in I_QUESTIONS:
                score = getattr(p.bia, f"{key}_score", None)
                row.append(_score_str(score, SCORE_BIA))
            row.append(_score_str(p.bia.integrity_score, SCORE_LABELS))
            _xl_append(ws, row)
        _xl_autowidth(ws)

    if "vertrouwelijkheid" in sections:
        ws = wb.create_sheet("Vertrouwelijkheid (V)")
        q_headers = ["Proces"] + [label for _, label in V_QUESTIONS] + ["Eindscore V"]
        _xl_header(ws, q_headers)
        for p in processes:
            if p.bia is None:
                continue
            row = [f"{p.code} – {p.name}"]
            for key, _ in V_QUESTIONS:
                score = getattr(p.bia, f"{key}_score", None)
                row.append(_score_str(score, SCORE_BIA))
            row.append(_score_str(p.bia.confidentiality_score, SCORE_LABELS))
            _xl_append(ws, row)
        _xl_autowidth(ws)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


def _build_xlsx_business_context(db: Session, sections: list[str]) -> io.BytesIO:
    from openpyxl import Workbook
    processes = _get_processes(db)
    wb = Workbook()
    wb.remove(wb.active)

    if "canvas" in sections:
        ws = wb.create_sheet("Canvas Blokken")
        _xl_header(ws, ["Proces", "Kern partners", "Kern activiteiten", "Waarde propositie",
                         "Klant relaties", "Klant segmenten", "Kanalen",
                         "Key resources", "Kostenstructuur", "Inkomstenstromen"])
        for p in processes:
            bc = p.business_context
            _xl_append(ws, [
                f"{p.code} – {p.name}",
                bc.key_partners or "" if bc else "",
                bc.key_activities or "" if bc else "",
                bc.value_proposition or "" if bc else "",
                bc.customer_relationships or "" if bc else "",
                bc.customer_segments or "" if bc else "",
                bc.channels or "" if bc else "",
                bc.key_resources or "" if bc else "",
                bc.cost_structure or "" if bc else "",
                bc.revenue_streams or "" if bc else "",
            ])
        _xl_autowidth(ws)

    if "wettelijk" in sections:
        ws = wb.create_sheet("Wettelijk & Stakeholders")
        _xl_header(ws, ["Proces", "Wettelijke basis", "Stakeholders", "Ketenpositie", "Key aspects"])
        for p in processes:
            bc = p.business_context
            _xl_append(ws, [
                f"{p.code} – {p.name}",
                bc.legal_basis or "" if bc else "",
                bc.stakeholders or "" if bc else "",
                bc.chain_position or "" if bc else "",
                bc.key_aspects or "" if bc else "",
            ])
        _xl_autowidth(ws)

    if "privacy" in sections:
        ws = wb.create_sheet("Privacy")
        _xl_header(ws, ["Proces", "Persoonsgegevens", "Bijzondere persoonsgegevens"])
        for p in processes:
            bc = p.business_context
            _xl_append(ws, [
                f"{p.code} – {p.name}",
                "Ja" if bc and bc.personal_data else "Nee",
                "Ja" if bc and bc.special_personal_data else "Nee",
            ])
        _xl_autowidth(ws)

    if "continuiteit" in sections:
        ws = wb.create_sheet("Continuïteit")
        _xl_header(ws, ["Proces", "Continuïteitsvereisten", "Review datum", "Notities"])
        for p in processes:
            bc = p.business_context
            _xl_append(ws, [
                f"{p.code} – {p.name}",
                bc.continuity_requirements or "" if bc else "",
                str(bc.review_date) if bc and bc.review_date else "",
                bc.notes or "" if bc else "",
            ])
        _xl_autowidth(ws)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


def _build_xlsx_ketenarchitectuur(db: Session, sections: list[str]) -> io.BytesIO:
    from openpyxl import Workbook
    processes = _get_processes(db)
    applications = _get_applications(db)
    wb = Workbook()
    wb.remove(wb.active)

    if "processen" in sections:
        ws = wb.create_sheet("Processen")
        _xl_header(ws, ["Code", "Naam", "Eigenaar", "Afdeling", "Kritiek", "# Applicaties", "B-score", "I-score", "V-score"])
        for p in processes:
            _xl_append(ws, [
                p.code, p.name, p.owner or "", p.department or "",
                "Ja" if p.is_critical else "Nee", len(p.applications),
                _score_str(p.bia.availability_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.integrity_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.confidentiality_score if p.bia else None, SCORE_LABELS),
            ])
        _xl_autowidth(ws)

    if "applicaties" in sections:
        ws = wb.create_sheet("Applicaties")
        _xl_header(ws, ["Code", "Naam", "Business owner", "Technisch owner", "# Processen"])
        for a in applications:
            _xl_append(ws, [a.code, a.name, a.business_owner or "", a.technical_owner or "", len(a.processes)])
        _xl_autowidth(ws)

    if "koppelingen" in sections:
        ws = wb.create_sheet("Koppelingen")
        _xl_header(ws, ["Proces code", "Proces naam", "Applicatie code", "Applicatie naam"])
        for p in processes:
            for a in p.applications:
                _xl_append(ws, [p.code, p.name, a.code, a.name])
        _xl_autowidth(ws)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


# ══════════════════════════════════════════════════════════════════════════════
#  DOCX BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def _build_docx_dashboard(db: Session, sections: list[str]) -> io.BytesIO:
    from docx import Document
    doc = Document()
    doc.add_heading("ProcesCheck – Dashboard Export", 0)
    processes = _get_processes(db)
    total = len(processes)

    def _cov(done):
        return round(done / total * 100) if total > 0 else 0

    if "kpi" in sections:
        doc.add_heading("KPI Overzicht", 1)
        complete = sum(1 for p in processes if _check_completeness(p)[0])
        bia_cov  = sum(1 for p in processes if p.bia is not None)
        bc_cov   = sum(1 for p in processes if p.business_context is not None)
        _doc_table(doc, ["Indicator", "Waarde", "%"], [
            ["Totaal processen", total, ""],
            ["Kritieke processen", sum(1 for p in processes if p.is_critical), f"{_cov(sum(1 for p in processes if p.is_critical))}%"],
            ["Volledig compleet", complete, f"{_cov(complete)}%"],
            ["BIA / BIV gedekt", bia_cov, f"{_cov(bia_cov)}%"],
            ["Business Context gedekt", bc_cov, f"{_cov(bc_cov)}%"],
        ])

    if "biv_verdeling" in sections:
        doc.add_heading("BIV-verdeling", 1)
        rows = []
        for dim, attr in [("Beschikbaarheid", "availability_score"), ("Integriteit", "integrity_score"), ("Vertrouwelijkheid", "confidentiality_score")]:
            counts = {1: 0, 2: 0, 3: 0, 4: 0, 5: 0}
            na = 0
            for p in processes:
                if p.bia is None:
                    na += 1
                else:
                    s = getattr(p.bia, attr, None)
                    if s in counts:
                        counts[s] += 1
                    else:
                        na += 1
            rows.append([dim, counts[1], counts[2], counts[3], counts[4], counts[5], na])
        _doc_table(doc, ["Dimensie", "Vitaal", "Hoog", "Midden", "Laag", "Minimaal", "N/A"], rows)

    if "kritieke_processen" in sections:
        doc.add_heading("Kritieke Processen", 1)
        rows = []
        for p in processes:
            if not p.is_critical:
                continue
            rows.append([
                p.code, p.name,
                _score_str(p.bia.availability_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.integrity_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.confidentiality_score if p.bia else None, SCORE_LABELS),
            ])
        _doc_table(doc, ["Code", "Naam", "B-score", "I-score", "V-score"], rows)

    if "review" in sections:
        doc.add_heading("Review Monitoring", 1)
        rev = _get_review_data(db)
        rows = [[cat, f"{done}/{total_n}", f"{pct}%"] for cat, (done, total_n, pct) in rev.items()]
        _doc_table(doc, ["Categorie", "Op tijd / Totaal", "Percentage"], rows)

    if "acties" in sections:
        doc.add_heading("Prioritaire Acties", 1)
        rows = [[a["code"], a["naam"], a["prioriteit"], a["reden"], a["ontbrekend"]]
                for a in _get_priority_actions(processes)]
        _doc_table(doc, ["Code", "Naam", "Prioriteit", "Reden", "Ontbrekende velden"], rows)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _build_docx_processes(db: Session, sections: list[str]) -> io.BytesIO:
    from docx import Document
    doc = Document()
    doc.add_heading("ProcesCheck – Processen Export", 0)
    processes = _get_processes(db)

    if "basis" in sections or "details" in sections:
        doc.add_heading("Processenlijst", 1)
        headers = []
        if "basis" in sections:
            headers += ["Code", "Naam", "Eigenaar", "Afdeling", "Kritiek"]
        if "details" in sections:
            headers += ["Beschrijving", "Doelstelling"]
        if "biv" in sections:
            headers += ["B", "I", "V"]
        rows = []
        for p in processes:
            row = []
            if "basis" in sections:
                row += [p.code, p.name, p.owner or "", p.department or "", "Ja" if p.is_critical else "Nee"]
            if "details" in sections:
                row += [p.description or "", p.objective or ""]
            if "biv" in sections:
                row += [
                    _score_str(p.bia.availability_score if p.bia else None, SCORE_LABELS),
                    _score_str(p.bia.integrity_score if p.bia else None, SCORE_LABELS),
                    _score_str(p.bia.confidentiality_score if p.bia else None, SCORE_LABELS),
                ]
            rows.append(row)
        _doc_table(doc, headers, rows)

    if "rto_rpo" in sections:
        doc.add_heading("RTO / RPO per Proces", 1)
        rows = []
        for p in processes:
            rows.append([p.code, p.name, _rto_str(p) or "N/A", _rpo_str(p) or "N/A"])
        _doc_table(doc, ["Code", "Naam", "RTO", "RPO"], rows)

    if "applicaties" in sections:
        doc.add_heading("Gekoppelde Applicaties per Proces", 1)
        rows = []
        for p in processes:
            rows.append([p.code, p.name, str(len(p.applications)), ", ".join(a.name for a in p.applications)])
        _doc_table(doc, ["Code", "Naam", "Aantal", "Applicaties"], rows)

    if "datums" in sections:
        doc.add_heading("Datums", 1)
        rows = [[p.code, p.name, str(p.last_assessment_date) if p.last_assessment_date else ""] for p in processes]
        _doc_table(doc, ["Code", "Naam", "Laatste beoordeling"], rows)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _build_docx_applications(db: Session, sections: list[str]) -> io.BytesIO:
    from docx import Document
    doc = Document()
    doc.add_heading("ProcesCheck – Applicaties Export", 0)
    applications = _get_applications(db)

    if "basis" in sections:
        doc.add_heading("Applicatielijst", 1)
        rows = [[a.code, a.name, a.business_owner or "", a.technical_owner or ""] for a in applications]
        _doc_table(doc, ["Code", "Naam", "Business owner", "Technisch owner"], rows)

    if "details" in sections:
        doc.add_heading("Details", 1)
        rows = [[a.code, a.name, a.description or "", a.notes or ""] for a in applications]
        _doc_table(doc, ["Code", "Naam", "Beschrijving", "Notities"], rows)

    if "review" in sections:
        doc.add_heading("Review Datum", 1)
        rows = [[a.code, a.name, str(a.review_date) if a.review_date else ""] for a in applications]
        _doc_table(doc, ["Code", "Naam", "Review datum"], rows)

    if "processen" in sections:
        doc.add_heading("Gekoppelde Processen per Applicatie", 1)
        rows = [[a.code, a.name, str(len(a.processes)), ", ".join(p.name for p in a.processes)] for a in applications]
        _doc_table(doc, ["Code", "Naam", "Aantal", "Processen"], rows)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _build_docx_bia(db: Session, sections: list[str]) -> io.BytesIO:
    from docx import Document
    doc = Document()
    doc.add_heading("ProcesCheck – BIA & BIV-Classificatie Export", 0)
    processes = _get_processes(db)

    if "eindscores" in sections:
        doc.add_heading("Eindscores BIV per Proces", 1)
        rows = []
        for p in processes:
            rows.append([
                p.code, p.name,
                _score_str(p.bia.availability_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.integrity_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.confidentiality_score if p.bia else None, SCORE_LABELS),
            ])
        _doc_table(doc, ["Code", "Naam", "B-score", "I-score", "V-score"], rows)

    if "algemeen" in sections:
        doc.add_heading("Algemene BIA-informatie", 1)
        for p in processes:
            if p.bia is None:
                continue
            doc.add_heading(f"{p.code} – {p.name}", 2)
            _doc_kv(doc, [
                ("Interviewer", p.bia.interviewer_name or ""),
                ("Interview datum", str(p.bia.interview_date) if p.bia.interview_date else ""),
                ("Beschrijving", p.bia.general_description or ""),
                ("Ketenafhankelijkheden", p.bia.chain_dependencies or ""),
                ("Afwijking motivatie (eigenaar)", p.bia.owner_deviation_motivation or ""),
                ("Notities", p.bia.notes or ""),
            ])

    if "beschikbaarheid" in sections:
        doc.add_heading("Beschikbaarheid (B) – Antwoorden per Proces", 1)
        headers = ["Proces"] + [label for _, label in B_QUESTIONS] + ["Eindscore B"]
        rows = []
        for p in processes:
            if p.bia is None:
                continue
            row = [f"{p.code} – {p.name}"]
            for key, _ in B_QUESTIONS:
                score = getattr(p.bia, f"{key}_score", None)
                row.append(_score_str(score, SCORE_BIA))
            row.append(_score_str(p.bia.availability_score, SCORE_LABELS))
            rows.append(row)
        _doc_table(doc, headers, rows)

    if "integriteit" in sections:
        doc.add_heading("Integriteit (I) – Antwoorden per Proces", 1)
        headers = ["Proces"] + [label for _, label in I_QUESTIONS] + ["Eindscore I"]
        rows = []
        for p in processes:
            if p.bia is None:
                continue
            row = [f"{p.code} – {p.name}"]
            for key, _ in I_QUESTIONS:
                score = getattr(p.bia, f"{key}_score", None)
                row.append(_score_str(score, SCORE_BIA))
            row.append(_score_str(p.bia.integrity_score, SCORE_LABELS))
            rows.append(row)
        _doc_table(doc, headers, rows)

    if "vertrouwelijkheid" in sections:
        doc.add_heading("Vertrouwelijkheid (V) – Antwoorden per Proces", 1)
        headers = ["Proces"] + [label for _, label in V_QUESTIONS] + ["Eindscore V"]
        rows = []
        for p in processes:
            if p.bia is None:
                continue
            row = [f"{p.code} – {p.name}"]
            for key, _ in V_QUESTIONS:
                score = getattr(p.bia, f"{key}_score", None)
                row.append(_score_str(score, SCORE_BIA))
            row.append(_score_str(p.bia.confidentiality_score, SCORE_LABELS))
            rows.append(row)
        _doc_table(doc, headers, rows)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _build_docx_business_context(db: Session, sections: list[str]) -> io.BytesIO:
    from docx import Document
    doc = Document()
    doc.add_heading("ProcesCheck – Procescontext Export", 0)
    processes = _get_processes(db)

    if "canvas" in sections:
        doc.add_heading("Canvas Blokken", 1)
        canvas_fields = [
            ("key_partners", "Kern partners"),
            ("key_activities", "Kern activiteiten"),
            ("value_proposition", "Waarde propositie"),
            ("customer_relationships", "Klant relaties"),
            ("customer_segments", "Klant segmenten"),
            ("channels", "Kanalen"),
            ("key_resources", "Key resources"),
            ("cost_structure", "Kostenstructuur"),
            ("revenue_streams", "Inkomstenstromen"),
        ]
        for p in processes:
            bc = p.business_context
            doc.add_heading(f"{p.code} – {p.name}", 2)
            _doc_kv(doc, [(label, getattr(bc, field, "") or "" if bc else "") for field, label in canvas_fields])

    if "wettelijk" in sections:
        doc.add_heading("Wettelijk & Stakeholders", 1)
        rows = []
        for p in processes:
            bc = p.business_context
            rows.append([
                f"{p.code} – {p.name}",
                bc.legal_basis or "" if bc else "",
                bc.stakeholders or "" if bc else "",
                bc.chain_position or "" if bc else "",
            ])
        _doc_table(doc, ["Proces", "Wettelijke basis", "Stakeholders", "Ketenpositie"], rows)

    if "privacy" in sections:
        doc.add_heading("Privacy", 1)
        rows = []
        for p in processes:
            bc = p.business_context
            rows.append([
                f"{p.code} – {p.name}",
                "Ja" if bc and bc.personal_data else "Nee",
                "Ja" if bc and bc.special_personal_data else "Nee",
            ])
        _doc_table(doc, ["Proces", "Persoonsgegevens", "Bijzondere persoonsgegevens"], rows)

    if "continuiteit" in sections:
        doc.add_heading("Continuïteit", 1)
        rows = []
        for p in processes:
            bc = p.business_context
            rows.append([
                f"{p.code} – {p.name}",
                bc.continuity_requirements or "" if bc else "",
                str(bc.review_date) if bc and bc.review_date else "",
            ])
        _doc_table(doc, ["Proces", "Continuïteitsvereisten", "Review datum"], rows)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _build_docx_ketenarchitectuur(db: Session, sections: list[str]) -> io.BytesIO:
    from docx import Document
    doc = Document()
    doc.add_heading("ProcesCheck – Ketenarchitectuur Export", 0)
    processes = _get_processes(db)
    applications = _get_applications(db)

    if "processen" in sections:
        doc.add_heading("Processen", 1)
        rows = [[p.code, p.name, p.owner or "", p.department or "", "Ja" if p.is_critical else "Nee"] for p in processes]
        _doc_table(doc, ["Code", "Naam", "Eigenaar", "Afdeling", "Kritiek"], rows)

    if "applicaties" in sections:
        doc.add_heading("Applicaties", 1)
        rows = [[a.code, a.name, a.business_owner or "", a.technical_owner or ""] for a in applications]
        _doc_table(doc, ["Code", "Naam", "Business owner", "Technisch owner"], rows)

    if "koppelingen" in sections:
        doc.add_heading("Koppelingen Proces – Applicatie", 1)
        rows = []
        for p in processes:
            for a in p.applications:
                rows.append([p.code, p.name, a.code, a.name])
        _doc_table(doc, ["Proces code", "Proces naam", "Applicatie code", "Applicatie naam"], rows)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# ══════════════════════════════════════════════════════════════════════════════
#  PROCESDOSSIER (één proces, compleet)
# ══════════════════════════════════════════════════════════════════════════════

def _build_docx_process_dossier(p: Process) -> io.BytesIO:
    from docx import Document
    doc = Document()
    doc.add_heading(f"Procesdossier – {p.code} {p.name}", 0)
    doc.add_paragraph(f"Gegenereerd op {_date.today().strftime('%d-%m-%Y')} vanuit ProcesCheck.")

    # 1. Basisgegevens
    doc.add_heading("1. Basisgegevens", 1)
    basis = [
        ("Procescode", p.code),
        ("Naam", p.name),
        ("Eigenaar", p.owner or ""),
        ("Afdeling", p.department or ""),
        ("Kritiek proces", "Ja" if p.is_critical else "Nee"),
    ]
    if p.is_critical:
        basis.append(("Reden kritiek", p.critical_reason or ""))
    basis += [
        ("Beschrijving", p.description or ""),
        ("Doelstelling", p.objective or ""),
        ("Laatste beoordeling", str(p.last_assessment_date) if p.last_assessment_date else ""),
        ("Notities", p.notes or ""),
    ]
    _doc_kv(doc, basis)

    # 2. BIV-classificatie en continuïteitsparameters
    doc.add_heading("2. BIV-classificatie & continuïteitsparameters", 1)
    if p.bia is None:
        doc.add_paragraph("Nog geen BIA ingevuld voor dit proces.")
    else:
        _doc_kv(doc, [
            ("Beschikbaarheid (B)", _score_str(p.bia.availability_score, SCORE_LABELS)),
            ("Integriteit (I)", _score_str(p.bia.integrity_score, SCORE_LABELS)),
            ("Vertrouwelijkheid (V)", _score_str(p.bia.confidentiality_score, SCORE_LABELS)),
            ("RTO (uit B1)", _rto_str(p)),
            ("RPO (uit B2)", _rpo_str(p)),
        ])

    # 3. BIA – algemene informatie en antwoorden met argumentatie
    doc.add_heading("3. Business Impact Analyse", 1)
    if p.bia is None:
        doc.add_paragraph("Nog geen BIA ingevuld voor dit proces.")
    else:
        _doc_kv(doc, [
            ("Interviewer", p.bia.interviewer_name or ""),
            ("Interview datum", str(p.bia.interview_date) if p.bia.interview_date else ""),
            ("Algemene beschrijving", p.bia.general_description or ""),
            ("Ketenafhankelijkheden", p.bia.chain_dependencies or ""),
            ("Afwijking motivatie (eigenaar)", p.bia.owner_deviation_motivation or ""),
            ("Notities", p.bia.notes or ""),
        ])
        doc.add_heading("Antwoorden per vraag", 2)
        rows = []
        for key, label in B_QUESTIONS + I_QUESTIONS + V_QUESTIONS:
            score = getattr(p.bia, f"{key}_score", None)
            arg = getattr(p.bia, f"{key}_arg", None)
            rows.append([label, _score_str(score, SCORE_BIA), arg or ""])
        _doc_table(doc, ["Vraag", "Antwoord", "Argumentatie"], rows)

    # 4. Procescontext
    doc.add_heading("4. Procescontext", 1)
    bc = p.business_context
    if bc is None:
        doc.add_paragraph("Nog geen procescontext ingevuld voor dit proces.")
    else:
        _doc_kv(doc, [
            ("Kernpartners", bc.key_partners or ""),
            ("Kernactiviteiten", bc.key_activities or ""),
            ("Kernresources", bc.key_resources or ""),
            ("Waardepropositie", bc.value_proposition or ""),
            ("Klantrelaties", bc.customer_relationships or ""),
            ("Kanalen", bc.channels or ""),
            ("Klantsegmenten", bc.customer_segments or ""),
            ("Kostenstructuur", bc.cost_structure or ""),
            ("Inkomstenstromen", bc.revenue_streams or ""),
            ("Wettelijke basis", bc.legal_basis or ""),
            ("Stakeholders", bc.stakeholders or ""),
            ("Ketenpositie", bc.chain_position or ""),
            ("Continuïteitsvereisten", bc.continuity_requirements or ""),
            ("Key aspects", bc.key_aspects or ""),
            ("Persoonsgegevens", "Ja" if bc.personal_data else "Nee"),
            ("Bijzondere persoonsgegevens", "Ja" if bc.special_personal_data else "Nee"),
            ("Review datum", str(bc.review_date) if bc.review_date else ""),
            ("Notities", bc.notes or ""),
        ])

    # 5. Gekoppelde applicaties
    doc.add_heading("5. Gekoppelde applicaties", 1)
    if not p.applications:
        doc.add_paragraph("Geen applicaties gekoppeld aan dit proces.")
    else:
        rows = [
            [a.code, a.name, a.business_owner or "", a.technical_owner or "", a.description or ""]
            for a in p.applications
        ]
        _doc_table(doc, ["Code", "Naam", "Business owner", "Technisch owner", "Beschrijving"], rows)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# ══════════════════════════════════════════════════════════════════════════════
#  PPTX BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def _build_pptx_dashboard(db: Session, sections: list[str]) -> io.BytesIO:
    from pptx import Presentation
    prs = Presentation()
    _prs_title_slide(prs, "ProcesCheck – Dashboard", f"Export {_date.today()}")
    processes = _get_processes(db)
    total = len(processes)

    def _cov(done):
        return round(done / total * 100) if total > 0 else 0

    if "kpi" in sections:
        complete = sum(1 for p in processes if _check_completeness(p)[0])
        bia_cov  = sum(1 for p in processes if p.bia is not None)
        bc_cov   = sum(1 for p in processes if p.business_context is not None)
        _prs_table_slide(prs, "KPI Overzicht", ["Indicator", "Waarde", "%"], [
            ["Totaal processen", total, ""],
            ["Kritieke processen", sum(1 for p in processes if p.is_critical), f"{_cov(sum(1 for p in processes if p.is_critical))}%"],
            ["Volledig compleet", complete, f"{_cov(complete)}%"],
            ["BIA / BIV gedekt", bia_cov, f"{_cov(bia_cov)}%"],
            ["Business Context gedekt", bc_cov, f"{_cov(bc_cov)}%"],
        ])

    if "biv_verdeling" in sections:
        rows = []
        for dim, attr in [("Beschikbaarheid", "availability_score"), ("Integriteit", "integrity_score"), ("Vertrouwelijkheid", "confidentiality_score")]:
            counts = {1: 0, 2: 0, 3: 0, 4: 0, 5: 0}
            na = 0
            for p in processes:
                if p.bia is None:
                    na += 1
                else:
                    s = getattr(p.bia, attr, None)
                    if s in counts:
                        counts[s] += 1
                    else:
                        na += 1
            rows.append([dim, counts[1], counts[2], counts[3], counts[4], counts[5], na])
        _prs_table_slide(prs, "BIV-verdeling", ["Dimensie", "Vitaal", "Hoog", "Midden", "Laag", "Minimaal", "N/A"], rows)

    if "kritieke_processen" in sections:
        rows = []
        for p in processes:
            if not p.is_critical:
                continue
            rows.append([
                p.code, p.name,
                _score_str(p.bia.availability_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.integrity_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.confidentiality_score if p.bia else None, SCORE_LABELS),
            ])
        _prs_table_slide(prs, "Kritieke Processen", ["Code", "Naam", "B", "I", "V"], rows)

    if "review" in sections:
        rev = _get_review_data(db)
        rows = [[cat, f"{done}/{total_n}", f"{pct}%"] for cat, (done, total_n, pct) in rev.items()]
        _prs_table_slide(prs, "Review Monitoring", ["Categorie", "Op tijd / Totaal", "%"], rows)

    if "acties" in sections:
        rows = [[a["code"], a["naam"], a["prioriteit"], a["reden"]] for a in _get_priority_actions(processes)]
        _prs_table_slide(prs, "Prioritaire Acties", ["Code", "Naam", "Prioriteit", "Reden"], rows)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _build_pptx_processes(db: Session, sections: list[str]) -> io.BytesIO:
    from pptx import Presentation
    prs = Presentation()
    _prs_title_slide(prs, "ProcesCheck – Processen", f"Export {_date.today()}")
    processes = _get_processes(db)

    if "basis" in sections or "biv" in sections:
        headers = []
        if "basis" in sections: headers += ["Code", "Naam", "Eigenaar", "Afdeling", "Kritiek"]
        if "biv" in sections: headers += ["B", "I", "V"]
        rows = []
        for p in processes:
            row = []
            if "basis" in sections:
                row += [p.code, p.name, p.owner or "", p.department or "", "Ja" if p.is_critical else "Nee"]
            if "biv" in sections:
                row += [
                    _score_str(p.bia.availability_score if p.bia else None, SCORE_LABELS),
                    _score_str(p.bia.integrity_score if p.bia else None, SCORE_LABELS),
                    _score_str(p.bia.confidentiality_score if p.bia else None, SCORE_LABELS),
                ]
            rows.append(row)
        _prs_table_slide(prs, "Processenlijst", headers, rows)

    if "rto_rpo" in sections:
        rows = []
        for p in processes:
            rows.append([p.code, p.name, _rto_str(p) or "N/A", _rpo_str(p) or "N/A"])
        _prs_table_slide(prs, "RTO / RPO per Proces", ["Code", "Naam", "RTO", "RPO"], rows)

    if "applicaties" in sections:
        rows = [[p.code, p.name, str(len(p.applications)), ", ".join(a.name for a in p.applications)] for p in processes]
        _prs_table_slide(prs, "Gekoppelde Applicaties", ["Code", "Naam", "#", "Applicaties"], rows)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _build_pptx_applications(db: Session, sections: list[str]) -> io.BytesIO:
    from pptx import Presentation
    prs = Presentation()
    _prs_title_slide(prs, "ProcesCheck – Applicaties", f"Export {_date.today()}")
    applications = _get_applications(db)

    if "basis" in sections:
        rows = [[a.code, a.name, a.business_owner or "", a.technical_owner or ""] for a in applications]
        _prs_table_slide(prs, "Applicatielijst", ["Code", "Naam", "Business owner", "Technisch owner"], rows)

    if "processen" in sections:
        rows = [[a.code, a.name, str(len(a.processes)), ", ".join(p.name for p in a.processes)] for a in applications]
        _prs_table_slide(prs, "Gekoppelde Processen", ["Code", "Naam", "#", "Processen"], rows)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _build_pptx_bia(db: Session, sections: list[str]) -> io.BytesIO:
    from pptx import Presentation
    prs = Presentation()
    _prs_title_slide(prs, "ProcesCheck – BIA & BIV-Classificatie", f"Export {_date.today()}")
    processes = _get_processes(db)

    if "eindscores" in sections:
        rows = []
        for p in processes:
            rows.append([
                p.code, p.name,
                _score_str(p.bia.availability_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.integrity_score if p.bia else None, SCORE_LABELS),
                _score_str(p.bia.confidentiality_score if p.bia else None, SCORE_LABELS),
            ])
        _prs_table_slide(prs, "Eindscores BIV", ["Code", "Naam", "B-score", "I-score", "V-score"], rows)

    if "beschikbaarheid" in sections:
        headers = ["Proces"] + [label for _, label in B_QUESTIONS] + ["Eindscore B"]
        rows = []
        for p in processes:
            if p.bia is None:
                continue
            row = [f"{p.code} – {p.name}"]
            for key, _ in B_QUESTIONS:
                score = getattr(p.bia, f"{key}_score", None)
                row.append(_score_str(score, SCORE_BIA) if score else "")
            row.append(_score_str(p.bia.availability_score, SCORE_LABELS))
            rows.append(row)
        _prs_table_slide(prs, "Beschikbaarheid (B)", headers, rows)

    if "integriteit" in sections:
        headers = ["Proces"] + [label for _, label in I_QUESTIONS] + ["Eindscore I"]
        rows = []
        for p in processes:
            if p.bia is None:
                continue
            row = [f"{p.code} – {p.name}"]
            for key, _ in I_QUESTIONS:
                score = getattr(p.bia, f"{key}_score", None)
                row.append(_score_str(score, SCORE_BIA) if score else "")
            row.append(_score_str(p.bia.integrity_score, SCORE_LABELS))
            rows.append(row)
        _prs_table_slide(prs, "Integriteit (I)", headers, rows)

    if "vertrouwelijkheid" in sections:
        headers = ["Proces"] + [label for _, label in V_QUESTIONS] + ["Eindscore V"]
        rows = []
        for p in processes:
            if p.bia is None:
                continue
            row = [f"{p.code} – {p.name}"]
            for key, _ in V_QUESTIONS:
                score = getattr(p.bia, f"{key}_score", None)
                row.append(_score_str(score, SCORE_BIA) if score else "")
            row.append(_score_str(p.bia.confidentiality_score, SCORE_LABELS))
            rows.append(row)
        _prs_table_slide(prs, "Vertrouwelijkheid (V)", headers, rows)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _build_pptx_business_context(db: Session, sections: list[str]) -> io.BytesIO:
    from pptx import Presentation
    prs = Presentation()
    _prs_title_slide(prs, "ProcesCheck – Procescontext", f"Export {_date.today()}")
    processes = _get_processes(db)

    if "canvas" in sections:
        rows = []
        for p in processes:
            bc = p.business_context
            rows.append([f"{p.code} – {p.name}",
                         bc.key_partners or "" if bc else "",
                         bc.key_activities or "" if bc else "",
                         bc.value_proposition or "" if bc else ""])
        _prs_table_slide(prs, "Canvas – Kernblokken", ["Proces", "Kern partners", "Kern activiteiten", "Waarde propositie"], rows)

    if "privacy" in sections:
        rows = []
        for p in processes:
            bc = p.business_context
            rows.append([f"{p.code} – {p.name}",
                         "Ja" if bc and bc.personal_data else "Nee",
                         "Ja" if bc and bc.special_personal_data else "Nee"])
        _prs_table_slide(prs, "Privacy per Proces", ["Proces", "Persoonsgegevens", "Bijzondere PG"], rows)

    if "wettelijk" in sections:
        rows = []
        for p in processes:
            bc = p.business_context
            rows.append([f"{p.code} – {p.name}", bc.legal_basis or "" if bc else "", bc.stakeholders or "" if bc else ""])
        _prs_table_slide(prs, "Wettelijk & Stakeholders", ["Proces", "Wettelijke basis", "Stakeholders"], rows)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _build_pptx_ketenarchitectuur(db: Session, sections: list[str]) -> io.BytesIO:
    from pptx import Presentation
    prs = Presentation()
    _prs_title_slide(prs, "ProcesCheck – Ketenarchitectuur", f"Export {_date.today()}")
    processes = _get_processes(db)
    applications = _get_applications(db)

    if "processen" in sections:
        rows = [[p.code, p.name, p.owner or "", p.department or "", "Ja" if p.is_critical else "Nee"] for p in processes]
        _prs_table_slide(prs, "Processen", ["Code", "Naam", "Eigenaar", "Afdeling", "Kritiek"], rows)

    if "applicaties" in sections:
        rows = [[a.code, a.name, a.business_owner or "", a.technical_owner or ""] for a in applications]
        _prs_table_slide(prs, "Applicaties", ["Code", "Naam", "Business owner", "Technisch owner"], rows)

    if "koppelingen" in sections:
        rows = []
        for p in processes:
            for a in p.applications:
                rows.append([p.code, p.name, a.code, a.name])
        _prs_table_slide(prs, "Koppelingen", ["Proces code", "Proces naam", "Applicatie code", "Applicatie naam"], rows)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ══════════════════════════════════════════════════════════════════════════════
#  ENDPOINTS
# ══════════════════════════════════════════════════════════════════════════════

XLSX_BUILDERS = {
    "dashboard":         _build_xlsx_dashboard,
    "processes":         _build_xlsx_processes,
    "applications":      _build_xlsx_applications,
    "bia":               _build_xlsx_bia,
    "business-context":  _build_xlsx_business_context,
    "ketenarchitectuur": _build_xlsx_ketenarchitectuur,
}

DOCX_BUILDERS = {
    "dashboard":         _build_docx_dashboard,
    "processes":         _build_docx_processes,
    "applications":      _build_docx_applications,
    "bia":               _build_docx_bia,
    "business-context":  _build_docx_business_context,
    "ketenarchitectuur": _build_docx_ketenarchitectuur,
}

PPTX_BUILDERS = {
    "dashboard":         _build_pptx_dashboard,
    "processes":         _build_pptx_processes,
    "applications":      _build_pptx_applications,
    "bia":               _build_pptx_bia,
    "business-context":  _build_pptx_business_context,
    "ketenarchitectuur": _build_pptx_ketenarchitectuur,
}

MODULE_NAMES = {
    "dashboard":         "dashboard",
    "processes":         "processen",
    "applications":      "applicaties",
    "bia":               "bia-biv",
    "business-context":  "procescontext",
    "ketenarchitectuur": "ketenarchitectuur",
}


@router.get("/process/{process_id}/docx")
def export_process_dossier(process_id: int, db: Session = Depends(get_db)):
    """Compleet dossier van één proces als Word-document."""
    p = db.get(Process, process_id)
    if not p:
        raise HTTPException(status_code=404, detail="Proces niet gevonden")
    buf = _build_docx_process_dossier(p)
    fname = f"procesdossier_{p.code}.docx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename={fname}"},
    )


@router.get("/{module}/xlsx")
def export_module_xlsx(module: str, sections: Optional[str] = None, db: Session = Depends(get_db)):
    if module not in XLSX_BUILDERS:
        raise HTTPException(status_code=404, detail=f"Module '{module}' niet gevonden.")
    selected = _parse_sections(sections, module)
    buf = XLSX_BUILDERS[module](db, selected)
    fname = f"procescheck_{MODULE_NAMES.get(module, module)}.xlsx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={fname}"},
    )


@router.get("/{module}/docx")
def export_module_docx(module: str, sections: Optional[str] = None, db: Session = Depends(get_db)):
    if module not in DOCX_BUILDERS:
        raise HTTPException(status_code=404, detail=f"Module '{module}' niet gevonden.")
    selected = _parse_sections(sections, module)
    buf = DOCX_BUILDERS[module](db, selected)
    fname = f"procescheck_{MODULE_NAMES.get(module, module)}.docx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename={fname}"},
    )


@router.get("/{module}/pptx")
def export_module_pptx(module: str, sections: Optional[str] = None, db: Session = Depends(get_db)):
    if module not in PPTX_BUILDERS:
        raise HTTPException(status_code=404, detail=f"Module '{module}' niet gevonden.")
    selected = _parse_sections(sections, module)
    buf = PPTX_BUILDERS[module](db, selected)
    fname = f"procescheck_{MODULE_NAMES.get(module, module)}.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={fname}"},
    )


# ── Backwards compatibility ───────────────────────────────────────────────────

@router.get("/processes.csv")
def export_processes_csv(db: Session = Depends(get_db)):
    processes = db.query(Process).order_by(Process.code).all()
    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerow([
        "Code", "Naam", "Beschrijving", "Doelstelling", "Eigenaar", "Afdeling",
        "Kritiek", "Reden kritiek", "Laatste beoordeling",
        "BIA B-score", "BIA I-score", "BIA V-score",
        "RTO", "RPO",
        "Aantal applicaties",
    ])
    for p in processes:
        writer.writerow([
            p.code, p.name, p.description or "", p.objective or "",
            p.owner or "", p.department or "",
            "Ja" if p.is_critical else "Nee",
            p.critical_reason or "",
            p.last_assessment_date or "",
            p.bia.availability_score if p.bia else "",
            p.bia.integrity_score if p.bia else "",
            p.bia.confidentiality_score if p.bia else "",
            _rto_str(p),
            _rpo_str(p),
            len(p.applications),
        ])
    output.seek(0)
    return StreamingResponse(
        iter([output.getvalue()]),
        media_type="text/csv",
        headers={"Content-Disposition": "attachment; filename=procescheck_export.csv"},
    )


@router.get("/processes.xlsx")
def export_processes_xlsx(db: Session = Depends(get_db)):
    return export_module_xlsx("processes", sections=None, db=db)
